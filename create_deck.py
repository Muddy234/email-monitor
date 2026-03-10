"""Generate a pitch deck PowerPoint walking through the email assistant pipeline."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand palette ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x1B, 0x1B, 0x2F)   # deep navy
BG_MID    = RGBColor(0x24, 0x24, 0x3E)   # card bg
ACCENT    = RGBColor(0x6C, 0x63, 0xFF)   # purple accent
ACCENT2   = RGBColor(0x00, 0xD2, 0xFF)   # cyan accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC0, 0xC0, 0xD0)   # body text
MUTED     = RGBColor(0x88, 0x88, 0xA0)   # subtle text
GREEN     = RGBColor(0x00, 0xE6, 0x96)
ORANGE    = RGBColor(0xFF, 0x9F, 0x43)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, default_size=16,
                     default_color=LIGHT, line_spacing=1.4):
    """lines: list of (text, {font overrides}) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(opts.get("size", default_size))
        p.font.color.rgb = opts.get("color", default_color)
        p.font.bold = opts.get("bold", False)
        p.font.name = opts.get("font", "Calibri")
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(opts.get("space_after", 4))
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    try:
        shape.text_frame.margin_top = Emu(0)
        shape.text_frame.margin_bottom = Emu(0)
    except Exception:
        pass
    return shape


def add_arrow(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 0.0
    return shape


def add_down_arrow(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def section_header(slide, number, title, subtitle=""):
    """Add a consistent section header with step number badge."""
    set_slide_bg(slide, BG_DARK)
    # Step badge
    add_rounded_rect(slide, Inches(0.7), Inches(0.5), Inches(1.2), Inches(0.5),
                     ACCENT, f"Step {number}", font_size=14, bold=True)
    # Title
    add_textbox(slide, Inches(2.1), Inches(0.4), Inches(9), Inches(0.6),
                title, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(2.1), Inches(1.0), Inches(9), Inches(0.5),
                    subtitle, font_size=16, color=MUTED)
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.7), Inches(1.55), Inches(11.9), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
            "Intelligent Email Assistant", font_size=44, bold=True, color=WHITE)
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            "Drafts Your Replies Before You Even Open Your Inbox",
            font_size=22, color=ACCENT2)

add_rich_textbox(slide, Inches(1), Inches(4.5), Inches(8), Inches(1.5), [
    ("An AI assistant that learns how you write and who you reply to.", {"size": 16, "color": LIGHT}),
    ("It reads your incoming mail, decides what needs a response, and writes a draft for you.",
     {"size": 16, "color": MUTED, "space_before": 8}),
])

# Decorative accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(1), Inches(4.2), Inches(3), Pt(3))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Overview (full pipeline at a glance)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7),
            "How It Works — The Big Picture", font_size=32, bold=True, color=WHITE)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.05), Inches(11.9), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Pipeline steps as cards
steps = [
    ("1", "Connect\nEmail",       "A browser extension\nlinks to your Outlook\nand mirrors your mail",   ACCENT),
    ("2", "Detect\nNew User",     "The system recognizes\na new user and waits\nfor enough email history", RGBColor(0x7C, 0x4D, 0xFF)),
    ("3", "Learn\nYour Style",    "Analyzes your past\nemails to understand\nhow you communicate",        RGBColor(0x9B, 0x59, 0xB6)),
    ("4", "Evaluate\nNew Mail",   "Each incoming email\nis scored and assessed:\ndoes it need a reply?",   RGBColor(0x00, 0xA8, 0xE8)),
    ("5", "Write\nthe Draft",     "AI composes a reply\nthat matches your tone\nand writing style",       ACCENT2),
    ("6", "Deliver\nto Outlook",  "The finished draft\nappears in your Outlook\nDrafts folder instantly", GREEN),
]

card_w = Inches(1.65)
card_h = Inches(3.8)
gap = Inches(0.35)
start_x = Inches(0.5)
card_y = Inches(1.5)
arrow_y = card_y + card_h / 2 - Inches(0.15)

for i, (num, title, desc, color) in enumerate(steps):
    x = start_x + i * (card_w + gap)

    # Card background
    add_rounded_rect(slide, x, card_y, card_w, card_h, BG_MID)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + card_w/2 - Inches(0.3),
                                    card_y + Inches(0.25), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, x + Inches(0.1), card_y + Inches(1.0), card_w - Inches(0.2), Inches(0.8),
                title, font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, x + Inches(0.1), card_y + Inches(1.85), card_w - Inches(0.2), Inches(1.8),
                desc, font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(steps) - 1:
        add_arrow(slide, x + card_w + Inches(0.05), arrow_y, Inches(0.25), Inches(0.3), color)

# Bottom note
add_textbox(slide, Inches(0.7), Inches(5.8), Inches(11), Inches(0.8),
            "Two AI tiers:  a fast model for analysis and triage  |  a premium model for writing drafts  |  batch processing cuts costs by 50%",
            font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Step 1: Email Sync
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 1, "Connect Your Email", "A lightweight browser extension links the system to your Outlook")

col1_x = Inches(0.7)
col2_x = Inches(6.8)
body_y = Inches(1.9)

add_rich_textbox(slide, col1_x, body_y, Inches(5.5), Inches(4.5), [
    ("How It Works", {"size": 20, "bold": True, "color": WHITE, "space_after": 10}),
    ("1.  User installs a small Chrome extension", {"size": 15, "color": LIGHT}),
    ("2.  The extension connects to Outlook Web automatically", {"size": 15, "color": LIGHT}),
    ("3.  Every 5 minutes, it copies new inbox and sent emails to our system", {"size": 15, "color": LIGHT}),
    ("4.  Incoming emails are queued for processing", {"size": 15, "color": LIGHT}),
    ("5.  Sent emails are kept for learning the user's writing style", {"size": 15, "color": LIGHT}),
])

# Right column — key details card
add_rounded_rect(slide, col2_x, body_y, Inches(5.5), Inches(3.5), BG_MID)
add_rich_textbox(slide, col2_x + Inches(0.3), body_y + Inches(0.2), Inches(4.9), Inches(3.0), [
    ("Key Details", {"size": 18, "bold": True, "color": ACCENT2, "space_after": 10}),
    ("Secure Connection", {"size": 14, "bold": True, "color": WHITE}),
    ("Uses Outlook's existing login session \u2014 we never see or store the user's password",
     {"size": 13, "color": LIGHT, "space_after": 8}),
    ("No Duplicates", {"size": 14, "bold": True, "color": WHITE}),
    ("Each email is identified uniquely \u2014 syncing the same email twice has no effect",
     {"size": 13, "color": LIGHT, "space_after": 8}),
    ("What We Capture", {"size": 14, "bold": True, "color": WHITE}),
    ("Sender, recipients, subject, body, conversation thread, and timestamps",
     {"size": 13, "color": LIGHT}),
])


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Step 2: Onboarding Trigger
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 2, "Detect New Users", "The system knows when it has enough data to start learning")

body_y = Inches(1.9)

add_rich_textbox(slide, Inches(0.7), body_y, Inches(5.5), Inches(4.5), [
    ("When Does Learning Begin?", {"size": 20, "bold": True, "color": WHITE, "space_after": 10}),
    ("All three must be true:", {"size": 15, "color": MUTED, "space_after": 6}),
    ("\u2713  The user hasn't been onboarded yet", {"size": 15, "color": GREEN}),
    ("\u2713  At least 20 emails have been synced", {"size": 15, "color": GREEN}),
    ("\u2713  No learning process is currently running", {"size": 15, "color": GREEN}),
    ("", {"size": 8}),
    ("The system checks for new users continuously in the background.",
     {"size": 14, "color": MUTED, "space_after": 6}),
    ("If anything goes wrong, it automatically retries.",
     {"size": 14, "color": MUTED}),
])

# Right side — flow diagram as stacked cards
add_rounded_rect(slide, Inches(7.0), body_y, Inches(5.0), Inches(1.0), BG_MID,
                 "System checks for new users", font_size=14)
add_down_arrow(slide, Inches(9.25), body_y + Inches(1.0), Inches(0.5), Inches(0.45), ACCENT)
add_rounded_rect(slide, Inches(7.0), body_y + Inches(1.5), Inches(5.0), Inches(1.0), BG_MID,
                 "Enough email history? Ready to learn?", font_size=14)
add_down_arrow(slide, Inches(9.25), body_y + Inches(2.5), Inches(0.5), Inches(0.45), ACCENT)
add_rounded_rect(slide, Inches(7.0), body_y + Inches(3.0), Inches(5.0), Inches(1.0), ACCENT,
                 "Begin the learning process \u2192", font_size=14, bold=True)


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Step 3: Onboarding Pipeline (overview)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, "Learning Your Style", "A 7-phase process that studies your email habits and writing voice")

phases = [
    ("Phase 1", "Gather\nEmails",     "Pulls 4 months of\nyour email history\nand filters out junk",          ACCENT),
    ("Phase 2", "Analyze\nPatterns",   "Studies who you\nreply to, how fast,\nand how often",                   RGBColor(0x7C, 0x4D, 0xFF)),
    ("Phase 3", "Read &\nExtract",     "AI reads a sample of\nyour emails to identify\ntopics and deadlines",   RGBColor(0x9B, 0x59, 0xB6)),
    ("Phase 4", "Build\nProfiles",     "AI maps your contacts\nand creates a personal\nwriting style guide",    RGBColor(0x00, 0xA8, 0xE8)),
    ("Phase 5", "Save\nResults",       "Stores everything the\nsystem learned about\nyou for future use",       ACCENT2),
    ("Phase 6", "Train\nthe Model",    "Builds a personalized\nscoring model to predict\nwhich emails matter",  GREEN),
    ("Phase 7", "Ready\nto Go",        "Onboarding complete\u2014\nthe system is now\nworking for you",         ORANGE),
]

card_w = Inches(1.55)
card_h = Inches(3.5)
gap = Inches(0.2)
start_x = Inches(0.4)
card_y = Inches(1.9)

for i, (phase, title, desc, color) in enumerate(phases):
    x = start_x + i * (card_w + gap)

    # Card
    add_rounded_rect(slide, x, card_y, card_w, card_h, BG_MID)

    # Phase label
    add_rounded_rect(slide, x + Inches(0.15), card_y + Inches(0.15),
                     card_w - Inches(0.3), Inches(0.35), color, phase, font_size=11, bold=True)

    # Title
    add_textbox(slide, x + Inches(0.08), card_y + Inches(0.6), card_w - Inches(0.16), Inches(0.5),
                title, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Desc
    add_textbox(slide, x + Inches(0.08), card_y + Inches(1.1), card_w - Inches(0.16), Inches(2.2),
                desc, font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Arrow
    if i < len(phases) - 1:
        add_arrow(slide, x + card_w + Inches(0.01), card_y + card_h/2 - Inches(0.12),
                  Inches(0.18), Inches(0.24), MUTED)

# Bottom note
add_textbox(slide, Inches(0.7), Inches(5.8), Inches(11.5), Inches(0.6),
            "Multiple phases run simultaneously for speed \u2014 the entire process completes in minutes",
            font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Step 3 Deep Dive: What Onboarding Learns
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, "What the System Learns", "A detailed look at the two types of analysis during onboarding")

body_y = Inches(1.9)

# Left column — Pattern analysis
add_rounded_rect(slide, Inches(0.7), body_y, Inches(5.6), Inches(4.8), BG_MID)
add_rich_textbox(slide, Inches(1.0), body_y + Inches(0.2), Inches(5.0), Inches(4.4), [
    ("Behavioral Pattern Analysis", {"size": 18, "bold": True, "color": ACCENT2, "space_after": 10}),
    ("Reply History", {"size": 14, "bold": True, "color": WHITE}),
    ("Which emails did you reply to? How quickly? This is the foundation for predictions.",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Contact Importance", {"size": 14, "bold": True, "color": WHITE}),
    ("How often you reply to each sender, your average response time, message volume",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Conversation Habits", {"size": 14, "bold": True, "color": WHITE}),
    ("How actively you participate in threads and whether you tend to start or join conversations",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Organization-Level Trends", {"size": 14, "bold": True, "color": WHITE}),
    ("Reply patterns grouped by company or domain (e.g., clients vs. internal team)",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Overall Profile", {"size": 14, "bold": True, "color": WHITE}),
    ("Your most active hours, the types of emails you send, and your general responsiveness",
     {"size": 12, "color": LIGHT}),
])

# Right column — AI understanding
add_rounded_rect(slide, Inches(6.8), body_y, Inches(5.6), Inches(4.8), BG_MID)
add_rich_textbox(slide, Inches(7.1), body_y + Inches(0.2), Inches(5.0), Inches(4.4), [
    ("AI-Powered Understanding", {"size": 18, "bold": True, "color": ACCENT, "space_after": 10}),
    ("Topic Identification", {"size": 14, "bold": True, "color": WHITE}),
    ("AI reads hundreds of your emails to learn what subjects and themes you deal with",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Writing Voice", {"size": 14, "bold": True, "color": WHITE}),
    ("Analyzes a sample of your sent emails to learn your greetings, sign-offs, tone, and phrasing",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Relationship Mapping", {"size": 14, "bold": True, "color": WHITE}),
    ("Figures out each contact's organization, role, and how important they are to you",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Topic Grouping", {"size": 14, "bold": True, "color": WHITE}),
    ("Organizes the topics you deal with into meaningful categories (e.g., legal, operations, finance)",
     {"size": 12, "color": LIGHT, "space_after": 8}),
    ("Personal Style Guide", {"size": 14, "bold": True, "color": WHITE}),
    ("Creates a guide for how you write to different people \u2014 more formal for outside counsel, casual for your team",
     {"size": 12, "color": LIGHT}),
])


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Step 3 Deep Dive: Scoring Model
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 3, "The Prioritization Model", "A personalized model that predicts which emails you'd actually reply to")

body_y = Inches(1.9)

add_rich_textbox(slide, Inches(0.7), body_y, Inches(5.5), Inches(4.8), [
    ("How Prioritization Works", {"size": 20, "bold": True, "color": WHITE, "space_after": 10}),
    ("The system builds a personalized model based on your", {"size": 14, "color": LIGHT, "space_after": 4}),
    ("past behavior. Each signal adds to or reduces the likelihood score.", {"size": 14, "color": LIGHT, "space_after": 12}),
    ("1.  Identify What Matters", {"size": 15, "color": ACCENT2}),
    ("     Which signals (sender, subject type, # of recipients) predict your replies?", {"size": 13, "color": LIGHT, "space_after": 6}),
    ("2.  Measure Each Signal's Weight", {"size": 15, "color": ACCENT2}),
    ("     How much does each factor increase or decrease your likelihood to respond?", {"size": 13, "color": LIGHT, "space_after": 6}),
    ("3.  Spot Recurring Emails", {"size": 15, "color": ACCENT2}),
    ("     Identifies newsletters, reports, and other emails that arrive on a regular schedule", {"size": 13, "color": LIGHT, "space_after": 6}),
    ("4.  Validate Against History", {"size": 15, "color": ACCENT2}),
    ("     Tests the model against your actual past replies to make sure it's accurate", {"size": 13, "color": LIGHT, "space_after": 6}),
    ("5.  Calibrate Confidence", {"size": 15, "color": ACCENT2}),
    ("     Converts raw scores into meaningful probabilities (e.g., \"80% likely to reply\")", {"size": 13, "color": LIGHT, "space_after": 6}),
    ("6.  Set Thresholds", {"size": 15, "color": ACCENT2}),
    ("     Decides what score is high enough to draft a reply vs. skip entirely", {"size": 13, "color": LIGHT}),
])

# Right side — output card
add_rounded_rect(slide, Inches(6.8), body_y, Inches(5.6), Inches(2.5), BG_MID)
add_rich_textbox(slide, Inches(7.1), body_y + Inches(0.2), Inches(5.0), Inches(2.1), [
    ("What the Model Produces", {"size": 18, "bold": True, "color": GREEN, "space_after": 10}),
    ("A complete scoring profile unique to each user:", {"size": 13, "color": MUTED, "space_after": 6}),
    ("\u2022  Weight for every signal (sender reputation, email type, etc.)", {"size": 14, "color": LIGHT}),
    ("\u2022  Probability calibration (so scores map to real-world likelihood)", {"size": 14, "color": LIGHT}),
    ("\u2022  Auto-skip threshold (below this = definitely ignore)", {"size": 14, "color": LIGHT}),
    ("\u2022  Draft threshold (above this = worth drafting a reply)", {"size": 14, "color": LIGHT}),
    ("\u2022  Known recurring senders (newsletters, reports, etc.)", {"size": 14, "color": LIGHT}),
])

add_rounded_rect(slide, Inches(6.8), body_y + Inches(2.8), Inches(5.6), Inches(1.5), BG_MID)
add_rich_textbox(slide, Inches(7.1), body_y + Inches(3.0), Inches(5.0), Inches(1.1), [
    ("Why This Approach?", {"size": 18, "bold": True, "color": ORANGE, "space_after": 8}),
    ("Transparent: you can always see why an email was prioritized.", {"size": 13, "color": LIGHT}),
    ("Works with small data \u2014 effective with as few as 20 emails.", {"size": 13, "color": LIGHT}),
    ("Lightweight \u2014 no expensive hardware required.", {"size": 13, "color": LIGHT}),
])


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Step 4: Runtime Email Processing
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 4, "Evaluating New Mail", "Every incoming email goes through a multi-stage evaluation")

body_y = Inches(1.9)

# Pipeline steps as vertical flow
steps = [
    ("Pick Up",   "New emails are claimed for processing \u2014 no email is ever handled twice",      ACCENT),
    ("Filter",    "Obvious noise is removed: auto-replies, out-of-office, newsletters, FYI messages", RGBColor(0x7C, 0x4D, 0xFF)),
    ("Score",     "The personalized model scores each email based on your past behavior",              RGBColor(0x00, 0xA8, 0xE8)),
    ("Enrich",    "Adds context: who is this sender? What's the conversation history?",                ACCENT2),
    ("Classify",  "AI makes the final call: does this email need a reply? What kind?",                 GREEN),
]

card_h = Inches(0.72)
gap = Inches(0.13)
for i, (label, desc, color) in enumerate(steps):
    y = body_y + i * (card_h + gap)
    add_rounded_rect(slide, Inches(0.7), y, Inches(1.4), card_h, color, label, font_size=15, bold=True)
    add_textbox(slide, Inches(2.3), y + Inches(0.15), Inches(5), card_h,
                desc, font_size=14, color=LIGHT)
    if i < len(steps) - 1:
        add_down_arrow(slide, Inches(1.2), y + card_h, Inches(0.4), Inches(0.12), MUTED)

# Right side — smart monitoring card
add_rounded_rect(slide, Inches(7.5), body_y, Inches(5.0), Inches(2.2), BG_MID)
add_rich_textbox(slide, Inches(7.8), body_y + Inches(0.15), Inches(4.4), Inches(1.9), [
    ("Always Watching", {"size": 18, "bold": True, "color": ACCENT2, "space_after": 8}),
    ("\u2022  Checks for new mail every 30 seconds", {"size": 13, "color": LIGHT}),
    ("\u2022  Ramps down when your inbox is quiet", {"size": 13, "color": LIGHT}),
    ("\u2022  Speeds back up the moment new mail arrives", {"size": 13, "color": LIGHT}),
    ("\u2022  Efficient \u2014 minimal system resources when idle", {"size": 13, "color": LIGHT}),
])

# Classification output card
add_rounded_rect(slide, Inches(7.5), body_y + Inches(2.5), Inches(5.0), Inches(2.2), BG_MID)
add_rich_textbox(slide, Inches(7.8), body_y + Inches(2.65), Inches(4.4), Inches(1.9), [
    ("The Verdict", {"size": 18, "bold": True, "color": GREEN, "space_after": 8}),
    ("\u2022  Does this email need a reply? (yes / no)", {"size": 13, "color": LIGHT}),
    ("\u2022  What kind of reply? (acknowledgment, answer, follow-up, etc.)", {"size": 13, "color": LIGHT}),
    ("\u2022  Why does it need a response?", {"size": 13, "color": LIGHT}),
    ("\u2022  How confident is the system in this decision?", {"size": 13, "color": LIGHT}),
    ("", {"size": 6}),
    ("Only recent emails that need a response move to drafting", {"size": 13, "color": ORANGE}),
])


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Step 5: Draft Generation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 5, "Writing the Draft", "AI composes a reply that sounds like you wrote it")

body_y = Inches(1.9)

# Prompt composition (left)
add_rounded_rect(slide, Inches(0.7), body_y, Inches(5.6), Inches(4.8), BG_MID)
add_rich_textbox(slide, Inches(1.0), body_y + Inches(0.2), Inches(5.0), Inches(4.4), [
    ("What the AI Considers", {"size": 18, "bold": True, "color": ACCENT2, "space_after": 12}),
    ("\u2022  The original email \u2014 who sent it, the subject, and the full message", {"size": 14, "color": LIGHT}),
    ("\u2022  Why this email needs a reply", {"size": 14, "color": LIGHT}),
    ("\u2022  What kind of reply is expected (e.g., confirmation, answer, follow-up)", {"size": 14, "color": LIGHT}),
    ("\u2022  Your relationship with the sender (client, colleague, vendor, etc.)", {"size": 14, "color": LIGHT}),
    ("\u2022  The conversation so far (prior messages in the thread)", {"size": 14, "color": LIGHT}),
    ("\u2022  How formal or casual to be with this person", {"size": 14, "color": LIGHT}),
    ("\u2022  Your personal writing style guide", {"size": 14, "color": LIGHT}),
    ("", {"size": 10}),
    ("The result: a draft that sounds like you \u2014 your greetings,", {"size": 13, "color": MUTED}),
    ("your sign-offs, your tone, your way of phrasing things.", {"size": 13, "color": MUTED}),
])

# Right — generation details
add_rich_textbox(slide, Inches(6.8), body_y, Inches(5.5), Inches(2.5), [
    ("Behind the Scenes", {"size": 20, "bold": True, "color": WHITE, "space_after": 10}),
    ("\u2022  Uses a premium AI model for high-quality writing", {"size": 14, "color": LIGHT}),
    ("\u2022  Each draft is validated before delivery (no blanks or errors)", {"size": 14, "color": LIGHT}),
    ("\u2022  Drafts are saved and queued for delivery to Outlook", {"size": 14, "color": LIGHT}),
    ("\u2022  If you've already edited a draft, the system won't overwrite your changes", {"size": 14, "color": LIGHT}),
])

add_rounded_rect(slide, Inches(6.8), body_y + Inches(2.8), Inches(5.5), Inches(2.0), BG_MID)
add_rich_textbox(slide, Inches(7.1), body_y + Inches(3.0), Inches(4.9), Inches(1.6), [
    ("Why a Premium Model for Drafts?", {"size": 18, "bold": True, "color": ACCENT, "space_after": 8}),
    ("This is the output you actually read and send \u2014", {"size": 13, "color": LIGHT}),
    ("quality matters most here.", {"size": 13, "color": LIGHT, "space_after": 6}),
    ("Batch processing keeps costs manageable even", {"size": 13, "color": LIGHT}),
    ("with the higher-quality model.", {"size": 13, "color": LIGHT}),
])


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Step 6: Draft Delivery
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, 6, "Deliver to Outlook", "The draft appears in your Drafts folder in real time")

body_y = Inches(1.9)

# Flow diagram — 4 stacked steps
flow_steps = [
    ("Live Connection", "The browser extension maintains a live connection to our system,\nlistening for new drafts as soon as they're ready", ACCENT),
    ("Draft Ready", "When a new draft is created, the extension is notified instantly\nand retrieves the original email's details (sender, subject)", RGBColor(0x00, 0xA8, 0xE8)),
    ("Create in Outlook", "The extension places the draft directly into your\nOutlook Drafts folder as a \"Re:\" reply to the original email", ACCENT2),
    ("Confirmed", "The draft is marked as delivered in our system\nYou'll find it waiting in your Drafts folder, ready to review", GREEN),
]

card_h = Inches(1.05)
gap = Inches(0.15)
for i, (title, desc, color) in enumerate(flow_steps):
    y = body_y + i * (card_h + gap)
    add_rounded_rect(slide, Inches(0.7), y, Inches(2.2), card_h, color, title, font_size=14, bold=True)
    add_textbox(slide, Inches(3.1), y + Inches(0.1), Inches(5), card_h,
                desc, font_size=13, color=LIGHT)
    if i < len(flow_steps) - 1:
        add_down_arrow(slide, Inches(1.6), y + card_h, Inches(0.4), Inches(0.14), MUTED)

# Right — user experience card
add_rounded_rect(slide, Inches(8.5), body_y, Inches(4.0), Inches(4.8), BG_MID)
add_rich_textbox(slide, Inches(8.8), body_y + Inches(0.2), Inches(3.4), Inches(4.4), [
    ("User Experience", {"size": 18, "bold": True, "color": GREEN, "space_after": 12}),
    ("The user simply opens", {"size": 15, "color": LIGHT}),
    ("their Outlook Drafts folder", {"size": 15, "color": LIGHT}),
    ("and finds ready-to-send", {"size": 15, "color": LIGHT}),
    ("replies waiting for them.", {"size": 15, "color": LIGHT, "space_after": 16}),
    ("They can:", {"size": 14, "color": MUTED, "space_after": 6}),
    ("\u2022  Send as-is", {"size": 14, "color": LIGHT}),
    ("\u2022  Edit and send", {"size": 14, "color": LIGHT}),
    ("\u2022  Discard", {"size": 14, "color": LIGHT, "space_after": 16}),
    ("Edits are tracked — the", {"size": 13, "color": MUTED}),
    ("system won't overwrite", {"size": 13, "color": MUTED}),
    ("user-modified drafts.", {"size": 13, "color": MUTED}),
])


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Architecture / Cost Summary
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "Design Principles & Cost Efficiency", font_size=32, bold=True, color=WHITE)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.1), Inches(11.9), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

body_y = Inches(1.5)

# Three cards
cards = [
    ("Smart AI Spending", [
        ("\u2022  Fast, affordable AI for analysis and triage", LIGHT),
        ("\u2022  Premium AI only where it matters \u2014 writing your drafts", LIGHT),
        ("\u2022  Batch processing cuts AI costs by 50%", LIGHT),
        ("\u2022  Scoring model runs locally \u2014 no expensive hardware", LIGHT),
    ], ACCENT),
    ("How Data Flows", [
        ("\u2022  Browser extension reads your Outlook mail", LIGHT),
        ("\u2022  Our system processes and scores each email", LIGHT),
        ("\u2022  Drafts are pushed back to the extension in real time", LIGHT),
        ("\u2022  Extension places the draft directly in Outlook", LIGHT),
    ], ACCENT2),
    ("Key Design Choices", [
        ("\u2022  Every email is processed exactly once \u2014 no duplicates", LIGHT),
        ("\u2022  System scales down when your inbox is quiet", LIGHT),
        ("\u2022  Each user gets their own personalized model", LIGHT),
        ("\u2022  Drafts match your writing voice, not a generic template", LIGHT),
    ], GREEN),
]

card_w = Inches(3.8)
gap = Inches(0.3)
for i, (title, items, color) in enumerate(cards):
    x = Inches(0.7) + i * (card_w + gap)
    add_rounded_rect(slide, x, body_y, card_w, Inches(4.2), BG_MID)
    add_rounded_rect(slide, x, body_y, card_w, Inches(0.5), color, title, font_size=15, bold=True)
    lines = [(title, {"size": 16, "bold": True, "color": color, "space_after": 10})]
    for text, col in items:
        lines.append((text, {"size": 13, "color": col, "space_after": 4}))
    add_rich_textbox(slide, x + Inches(0.2), body_y + Inches(0.7), card_w - Inches(0.4), Inches(3.3), lines)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — SWOT Analysis
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "SWOT Analysis", font_size=32, bold=True, color=WHITE)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.1), Inches(11.9), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

swot_y = Inches(1.4)
quad_w = Inches(5.85)
quad_h = Inches(2.8)
gap_x = Inches(0.6)
gap_y = Inches(0.3)

swot_data = [
    # Row 1
    [
        ("Strengths", GREEN, [
            "\u2022  Personalized per-user model \u2014 not one-size-fits-all",
            "\u2022  Drafts match the user's actual writing voice",
            "\u2022  Zero workflow change \u2014 drafts appear in Outlook natively",
            "\u2022  Cost-efficient: tiered AI models + batch processing (50% savings)",
            "\u2022  Fully autonomous after setup \u2014 no manual intervention",
        ]),
        ("Weaknesses", ORANGE, [
            "\u2022  Currently limited to Outlook Web + Chrome browser",
            "\u2022  Requires ~20 emails before the system can begin learning",
            "\u2022  Dependent on third-party AI providers for model access",
            "\u2022  Extension-based architecture limits mobile use cases",
            "\u2022  Draft quality is constrained by available email history",
        ]),
    ],
    # Row 2
    [
        ("Opportunities", ACCENT2, [
            "\u2022  Expand to Gmail, Apple Mail, and native desktop clients",
            "\u2022  Enterprise tier: team-wide intelligence, compliance controls",
            "\u2022  Adjacent features: meeting prep, task extraction, follow-up reminders",
            "\u2022  Rapidly growing market for AI-powered productivity tools",
            "\u2022  Mobile and native app integrations",
        ]),
        ("Threats", RGBColor(0xFF, 0x5C, 0x5C), [
            "\u2022  Microsoft Copilot and Google Gemini adding native email AI",
            "\u2022  Data privacy concerns may slow enterprise adoption",
            "\u2022  AI model pricing changes could impact unit economics",
            "\u2022  User trust barrier \u2014 hesitancy to let AI read personal email",
            "\u2022  Regulatory shifts around AI-generated communications",
        ]),
    ],
]

for row_i, row in enumerate(swot_data):
    for col_i, (title, color, bullets) in enumerate(row):
        x = Inches(0.7) + col_i * (quad_w + gap_x)
        y = swot_y + row_i * (quad_h + gap_y)

        # Card background
        add_rounded_rect(slide, x, y, quad_w, quad_h, BG_MID)

        # Colored header bar
        add_rounded_rect(slide, x, y, quad_w, Inches(0.5), color, title,
                         font_size=16, bold=True)

        # Bullet content
        bullet_lines = []
        for b in bullets:
            bullet_lines.append((b, {"size": 13, "color": LIGHT, "space_after": 3}))
        add_rich_textbox(slide, x + Inches(0.25), y + Inches(0.6),
                         quad_w - Inches(0.5), quad_h - Inches(0.7), bullet_lines)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Strengths Deep Dive (Internal)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_rounded_rect(slide, Inches(0.7), Inches(0.4), Inches(2.0), Inches(0.5),
                 GREEN, "STRENGTHS", font_size=14, bold=True)
add_textbox(slide, Inches(2.9), Inches(0.35), Inches(9), Inches(0.6),
            "What We Should Double Down On", font_size=28, bold=True, color=WHITE)
add_textbox(slide, Inches(2.9), Inches(0.9), Inches(9), Inches(0.4),
            "INTERNAL \u2014 How do we turn these into durable competitive advantages?",
            font_size=14, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.4), Inches(11.9), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = GREEN
line.line.fill.background()

body_y = Inches(1.7)

strengths_items = [
    ("Per-User Personalization",
     "This is our biggest moat. Generic AI email tools treat every user the same. Our per-user scoring model and style guide are hard to replicate without a similar architecture.",
     "How do we make the personalization even stickier? The longer someone uses it, the better it gets \u2014 can we quantify and communicate that flywheel to users?"),
    ("Zero Workflow Disruption",
     "Drafts land in the Outlook folder users already check. No new app, no new tab, no habit change required.",
     "This is a massive adoption advantage. We should lead with this in every pitch. Are there ways to make the experience even more invisible?"),
    ("Tiered AI + Batch Processing",
     "Using a fast model for triage and a premium model for drafts keeps quality high where it matters and costs low everywhere else. Batching saves 50%.",
     "As model prices drop, our margins improve automatically. Should we reinvest savings into quality (e.g., use the premium model more broadly) or pass savings to users?"),
    ("Fully Autonomous",
     "After the initial sync, users don't configure anything. Onboarding, model training, draft delivery \u2014 it all happens without intervention.",
     "Low-touch onboarding is great for scale. But does the \"black box\" feel make some users uneasy? Do we need a lightweight settings panel or a weekly summary email?"),
]

card_h = Inches(1.2)
gap = Inches(0.15)
for i, (title, desc, question) in enumerate(strengths_items):
    y = body_y + i * (card_h + gap)
    add_rounded_rect(slide, Inches(0.7), y, Inches(11.9), card_h, BG_MID)
    add_rich_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(11.3), card_h - Inches(0.2), [
        (title, {"size": 15, "bold": True, "color": GREEN, "space_after": 4}),
        (desc, {"size": 12, "color": LIGHT, "space_after": 4}),
        (f"\u25B8  {question}", {"size": 12, "color": ORANGE}),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — Weaknesses Deep Dive (Internal)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_rounded_rect(slide, Inches(0.7), Inches(0.4), Inches(2.3), Inches(0.5),
                 ORANGE, "WEAKNESSES", font_size=14, bold=True)
add_textbox(slide, Inches(3.2), Inches(0.35), Inches(9), Inches(0.6),
            "What We Need to Fix or Mitigate", font_size=28, bold=True, color=WHITE)
add_textbox(slide, Inches(3.2), Inches(0.9), Inches(9), Inches(0.4),
            "INTERNAL \u2014 Honest assessment of our current gaps and how to close them",
            font_size=14, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.4), Inches(11.9), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ORANGE
line.line.fill.background()

body_y = Inches(1.7)

weakness_items = [
    ("Outlook Web + Chrome Only",
     "We're locked into one email client on one browser. That's a small slice of the market \u2014 especially for mobile-heavy users and Gmail shops.",
     "Gmail support is the obvious next move. How much of the architecture is Outlook-specific vs. reusable? What's the realistic effort to abstract the email layer?"),
    ("Cold Start Problem (~20 emails)",
     "New users have to wait for enough email history before the system does anything useful. That's a gap between signup and first value.",
     "Can we offer a \"fast start\" mode using just the last few days of email? Or a manual config where users flag a few important senders up front?"),
    ("Third-Party AI Dependency",
     "We rely on Anthropic's API for all LLM work. A pricing change, rate limit, or outage directly impacts us.",
     "Should we build abstraction to swap providers? Or negotiate a committed-use agreement? What's our fallback if the API goes down for hours?"),
    ("No Mobile Experience",
     "The extension model doesn't work on phones. Users who triage email on mobile won't see our drafts until they open a desktop browser.",
     "Drafts sync to Outlook natively, so they do show up in the mobile Outlook app. Do we need to communicate this better, or build a companion mobile app?"),
]

card_h = Inches(1.2)
gap = Inches(0.15)
for i, (title, desc, question) in enumerate(weakness_items):
    y = body_y + i * (card_h + gap)
    add_rounded_rect(slide, Inches(0.7), y, Inches(11.9), card_h, BG_MID)
    add_rich_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(11.3), card_h - Inches(0.2), [
        (title, {"size": 15, "bold": True, "color": ORANGE, "space_after": 4}),
        (desc, {"size": 12, "color": LIGHT, "space_after": 4}),
        (f"\u25B8  {question}", {"size": 12, "color": ACCENT2}),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — Opportunities Deep Dive (Internal)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_rounded_rect(slide, Inches(0.7), Inches(0.4), Inches(2.8), Inches(0.5),
                 ACCENT2, "OPPORTUNITIES", font_size=14, bold=True)
add_textbox(slide, Inches(3.7), Inches(0.35), Inches(9), Inches(0.6),
            "Where We Can Grow", font_size=28, bold=True, color=WHITE)
add_textbox(slide, Inches(3.7), Inches(0.9), Inches(9), Inches(0.4),
            "INTERNAL \u2014 Which bets should we prioritize and in what order?",
            font_size=14, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.4), Inches(11.9), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT2
line.line.fill.background()

body_y = Inches(1.7)

opp_items = [
    ("Gmail + Multi-Platform Expansion",
     "Gmail is the largest email platform globally. Supporting it roughly doubles our addressable market overnight.",
     "What's the fastest path? Gmail API is well-documented. Can we reuse 80%+ of the worker pipeline and just swap the sync/delivery layer?"),
    ("Enterprise Tier",
     "Teams want shared context \u2014 org-wide contact intelligence, compliance guardrails, admin dashboards, and central billing.",
     "Enterprise sales cycles are longer but revenue per seat is much higher. What's the minimum viable enterprise feature set? Who do we pilot with?"),
    ("Adjacent Features",
     "We already understand the user's email deeply. That same intelligence could power meeting prep, action item extraction, follow-up reminders, and digest summaries.",
     "Which of these creates the most value with the least new infrastructure? Meeting prep and follow-up reminders seem closest to what we already do."),
    ("AI Productivity Market Tailwinds",
     "The market for AI productivity tools is growing rapidly. Businesses are actively looking for solutions, and willingness to pay is increasing.",
     "How do we position against horizontal tools (ChatGPT, Copilot) that do everything mediocrely vs. our deep vertical focus on email?"),
]

card_h = Inches(1.2)
gap = Inches(0.15)
for i, (title, desc, question) in enumerate(opp_items):
    y = body_y + i * (card_h + gap)
    add_rounded_rect(slide, Inches(0.7), y, Inches(11.9), card_h, BG_MID)
    add_rich_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(11.3), card_h - Inches(0.2), [
        (title, {"size": 15, "bold": True, "color": ACCENT2, "space_after": 4}),
        (desc, {"size": 12, "color": LIGHT, "space_after": 4}),
        (f"\u25B8  {question}", {"size": 12, "color": GREEN}),
    ])


# ══════════════════════════════════════════════════════════════
# SLIDE 16 — Threats Deep Dive (Internal)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

THREAT_RED = RGBColor(0xFF, 0x5C, 0x5C)

add_rounded_rect(slide, Inches(0.7), Inches(0.4), Inches(1.8), Inches(0.5),
                 THREAT_RED, "THREATS", font_size=14, bold=True)
add_textbox(slide, Inches(2.7), Inches(0.35), Inches(9), Inches(0.6),
            "What Could Hurt Us and How We Respond", font_size=28, bold=True, color=WHITE)
add_textbox(slide, Inches(2.7), Inches(0.9), Inches(9), Inches(0.4),
            "INTERNAL \u2014 Contingency thinking: what's our playbook if these materialize?",
            font_size=14, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(1.4), Inches(11.9), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = THREAT_RED
line.line.fill.background()

body_y = Inches(1.7)

threat_items = [
    ("Microsoft Copilot / Google Gemini",
     "The biggest platforms are building AI email features directly into their products. They have distribution we can't match and can ship to billions of users overnight.",
     "Their tools will be generic. Our advantage is deep personalization \u2014 per-user models, style matching, relationship context. Can we articulate that gap clearly in our messaging?"),
    ("Data Privacy & User Trust",
     "Letting an AI read your email is a big ask. Enterprise buyers will want SOC 2, data residency, and clear retention policies. Individual users may just feel uncomfortable.",
     "What's our data story? We should be proactive: publish a privacy whitepaper, pursue SOC 2 early, offer data deletion on demand. Transparency builds trust faster than reassurance."),
    ("AI Model Pricing Risk",
     "Our unit economics depend on current API pricing. A significant price increase from our AI provider would squeeze margins, especially at scale.",
     "Batch processing already cuts costs 50%. Should we explore open-source models as a fallback for the triage layer? That would reduce provider dependency for the cheaper tier."),
    ("Regulatory Uncertainty",
     "Governments are actively developing AI regulations. Rules around AI-generated communications, consent requirements, or disclosure mandates could impact our product.",
     "We should monitor EU AI Act and US state-level legislation. Can we build compliance features (e.g., AI-generated disclosure footers) proactively so we're ahead of mandates?"),
]

card_h = Inches(1.2)
gap = Inches(0.15)
for i, (title, desc, question) in enumerate(threat_items):
    y = body_y + i * (card_h + gap)
    add_rounded_rect(slide, Inches(0.7), y, Inches(11.9), card_h, BG_MID)
    add_rich_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(11.3), card_h - Inches(0.2), [
        (title, {"size": 15, "bold": True, "color": THREAT_RED, "space_after": 4}),
        (desc, {"size": 12, "color": LIGHT, "space_after": 4}),
        (f"\u25B8  {question}", {"size": 12, "color": ACCENT2}),
    ])


# ══════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════
output_path = r"C:\Users\NateMcBride\OneDrive - Arete Collective, L.P\Documents\Email_Monitor\Email_Assistant_Pipeline.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
